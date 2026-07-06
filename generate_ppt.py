from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def apply_light_theme(slide):
    # Very light gray/off-white background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(248, 250, 252) # Slate 50
    
    # Subtle top accent bar (Google Blue)
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.2))
    header.fill.solid()
    header.fill.fore_color.rgb = RGBColor(66, 133, 244) # Google Blue
    header.line.fill.background()
    
    # Very subtle bottom border
    glow = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.4), Inches(13.333), Inches(0.1))
    glow.fill.solid()
    glow.fill.fore_color.rgb = RGBColor(226, 232, 240) # Slate 200
    glow.line.fill.background()

def create_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255) # Pure White

    # Big Title
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(1.5))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "ChainGuardian"
    p.font.name = "Segoe UI"
    p.font.size = Pt(80)
    p.font.bold = True
    p.font.color.rgb = RGBColor(15, 23, 42) # Dark Slate
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(11.333), Inches(1))
    tf2 = txBox2.text_frame
    p2 = tf2.add_paragraph()
    p2.text = "AI-Powered Supply Chain Resilience\nGoogle APAC Hackathon"
    p2.font.name = "Segoe UI"
    p2.font.size = Pt(32)
    p2.font.color.rgb = RGBColor(100, 116, 139) # Slate 500
    p2.alignment = PP_ALIGN.CENTER

def create_card_slide(prs, title, description, cards):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_light_theme(slide)
    
    # Slide Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = title
    p.font.name = "Segoe UI"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(15, 23, 42)
    
    # Sub description
    txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(0.5))
    tf2 = txBox2.text_frame
    p2 = tf2.add_paragraph()
    p2.text = description
    p2.font.name = "Segoe UI"
    p2.font.size = Pt(24)
    p2.font.color.rgb = RGBColor(100, 116, 139)
    
    # Draw cards
    card_width = 3.6
    start_x = 0.6
    for i, card_text in enumerate(cards):
        x_pos = start_x + (i * (card_width + 0.4))
        
        # Card Background (Pure White)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x_pos), Inches(2.5), Inches(card_width), Inches(4))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(255, 255, 255)
        # Subtle border
        card.line.color.rgb = RGBColor(203, 213, 225) # Slate 300
        card.line.width = Pt(1.5)
        
        # Card Text
        ctxBox = slide.shapes.add_textbox(Inches(x_pos + 0.2), Inches(2.7), Inches(card_width - 0.4), Inches(3.6))
        ctf = ctxBox.text_frame
        ctf.word_wrap = True
        
        # Header (Number)
        cp1 = ctf.add_paragraph()
        cp1.text = f"0{i+1}"
        cp1.font.name = "Segoe UI"
        cp1.font.bold = True
        cp1.font.size = Pt(28)
        cp1.font.color.rgb = RGBColor(66, 133, 244) # Google Blue
        
        # Body
        cp2 = ctf.add_paragraph()
        cp2.text = "\n" + card_text
        cp2.font.name = "Segoe UI"
        cp2.font.size = Pt(20)
        cp2.font.color.rgb = RGBColor(51, 65, 85) # Slate 700

create_title(prs)

create_card_slide(prs, "The Vulnerability", "Why global supply chains fail during natural disasters:", [
    "REACTIVE SYSTEMS\nTraditional risk management relies on slow, manual geospatial data analysis.",
    "BLIND SPOTS\nDisasters like cyclones or earthquakes cause immediate port outages that are often realized too late.",
    "FINANCIAL LOSS\nDelayed routing decisions trap millions of dollars of cargo at sea."
])

create_card_slide(prs, "The Solution: ChainGuardian", "Predict and mitigate supply chain disruption in real-time.", [
    "LIVE INGESTION\nContinuously polls the Global Disaster Alert and Coordination System (GDACS).",
    "THREAT MAPPING\nSuperimposes live disaster blast radii over critical global shipping hubs.",
    "RAPIDS SPEED\nCalculates geospatial vulnerability matrices in milliseconds via NVIDIA RAPIDS."
])

create_card_slide(prs, "Gemini AI Copilot", "Data is useless without actionable intelligence.", [
    "AUTONOMOUS\nWhen a node is compromised, the AI automatically generates a mitigation strategy.",
    "FINANCIAL RISK\nInstantly computes 'Cargo at Risk' value vs 'Estimated Disruption Time'.",
    "REROUTE EXECUTION\nCopilot suggests specific fallback routes (e.g. 'Alpha Fallback via Panama Canal')."
])

output_path = os.path.join(os.getcwd(), "LIGHT_CORPORATE_PITCH.pptx")
prs.save(output_path)
print(f"Presentation saved to {output_path}")
